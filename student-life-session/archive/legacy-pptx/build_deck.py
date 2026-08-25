import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

prs = Presentation()
OLD_WIDTH = prs.slide_width  # 10in, the default 4:3 template width
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# The built-in template's master/layouts still carry 4:3-era placeholder
# geometry (e.g. a 9.0in-wide content box meant for a 10in-wide slide).
# Rescale every placeholder's left/width so they fill the new 16:9 canvas
# instead of leaving a large empty margin on the right.
RATIO = prs.slide_width / OLD_WIDTH

# Snapshot every placeholder's geometry BEFORE writing any of it back.
# Some layouts (Title and Content, Two Content, Comparison, Title Only)
# have no local <a:xfrm> at all and inherit their Title geometry straight
# from the master. If we read-then-write one placeholder at a time, the
# master gets rescaled first, so those inheriting layouts read the
# already-scaled master value and get double-scaled. Writing also forces
# python-pptx to materialize a brand-new local xfrm, which defaults any
# untouched dimension (top/height) to 0 instead of preserving the
# inherited value. Snapshotting all four dimensions for every placeholder
# up front, then writing them back explicitly, avoids both problems.
snapshots = []
for master in prs.slide_masters:
    for ph in master.placeholders:
        snapshots.append((ph, ph.left, ph.top, ph.width, ph.height))
    for layout in master.slide_layouts:
        for ph in layout.placeholders:
            snapshots.append((ph, ph.left, ph.top, ph.width, ph.height))

for ph, left, top, width, height in snapshots:
    ph.left = int(left * RATIO)
    ph.top = top
    ph.width = int(width * RATIO)
    ph.height = height

L_TITLE = prs.slide_layouts[0]
L_CONTENT = prs.slide_layouts[1]
L_SECTION = prs.slide_layouts[2]
L_TWO = prs.slide_layouts[3]
L_TITLE_ONLY = prs.slide_layouts[5]


def add_title_content_slide(title, layout=L_CONTENT):
    slide = prs.slides.add_slide(layout)
    slide.shapes.title.text = title
    body = None
    for ph in slide.placeholders:
        if ph.placeholder_format.idx == 1:
            body = ph
    return slide, body


def set_bullets(text_frame, items, font_size=20):
    text_frame.clear()
    text_frame.word_wrap = True
    for i, item in enumerate(items):
        p = text_frame.paragraphs[0] if i == 0 else text_frame.add_paragraph()
        if isinstance(item, tuple):
            text, level = item
        else:
            text, level = item, 0
        p.text = text
        p.level = level
        p.font.size = Pt(font_size)


def add_table(slide, rows, cols, left, top, width, height, data, col_widths=None, header=True):
    gframe = slide.shapes.add_table(rows, cols, left, top, width, height)
    table = gframe.table
    if col_widths:
        for i, w in enumerate(col_widths):
            table.columns[i].width = w
    for r, row_data in enumerate(data):
        for c, val in enumerate(row_data):
            cell = table.cell(r, c)
            cell.text = str(val)
            for p in cell.text_frame.paragraphs:
                p.font.size = Pt(14 if not (header and r == 0) else 15)
                if header and r == 0:
                    p.font.bold = True
    return table


def set_notes(slide, text):
    slide.notes_slide.notes_text_frame.text = text


# Slide 1 — Title
slide = prs.slides.add_slide(L_TITLE)
slide.shapes.title.text = "Student Life"
subtitle = slide.placeholders[1]
subtitle.text = "MScAC Orientation Week"
p2 = subtitle.text_frame.add_paragraph()
p2.text = "David — MScAC (Applied Math) 2025-2026"
p2.font.size = Pt(16)

# Slide 2 — Icebreaker pointer
slide, body = add_title_content_slide("Before we start... where in the GTA are you?", L_CONTENT)
set_bullets(body.text_frame, [
    "Scan the QR code or open the link to join the live map",
    "Duplicate a pin, drag it onto your neighbourhood",
], font_size=22)
qr_size = Inches(2.6)
slide.shapes.add_picture("icebreaker_qr.png", prs.slide_width - qr_size - Inches(0.5), Inches(1.6), width=qr_size, height=qr_size)
qr_caption = slide.shapes.add_textbox(prs.slide_width - qr_size - Inches(0.5), Inches(1.6) + qr_size + Inches(0.05), qr_size, Inches(0.4))
qr_caption.text_frame.text = "Live GTA map"
qr_caption.text_frame.paragraphs[0].font.size = Pt(12)
qr_caption.text_frame.paragraphs[0].font.color.rgb = RGBColor(0x77, 0x77, 0x77)

# Slide 4 — Section header
slide = prs.slides.add_slide(L_SECTION)
slide.shapes.title.text = "Weekends in Toronto"

# Slide 5 — Weekends: Summer/Fall
slide, body = add_title_content_slide("Weekends in Toronto: Summer/Fall", L_CONTENT)
set_bullets(body.text_frame, [
    "Parks & waterfront: High Park (cherry blossoms in spring), Toronto Islands (ferry ~$10 return), Waterfront Boardwalk",
    "Neighbourhoods: Chinatown, Koreatown (Bloor & Christie), Greektown (the Danforth), Little Italy (College St)",
    "Niagara Falls (day trip, GO train or bus or car) + Niagara-On-The-Lake",
    "Blue Jays game at Rogers Centre",
    "Canadian National Exhibition (late August)",
], font_size=18)

# Slide 6 — Weekends: Spring
slide, body = add_title_content_slide("Weekends in Toronto: Spring", L_CONTENT)
set_bullets(body.text_frame, [
    "Ice Skating at U of T Varsity Arena ($3.30 rental)",
    "Earl Bales Ski & Snowboard Center (North York)",
    "Blue Mountain Ski Resort",
    "Culture: Royal Ontario Museum (student rate, also free on Tuesdays)",
    "Cards to get: PRESTO, Toronto Public Libraries",
], font_size=18)
set_notes(slide, "This is a starting list, not exhaustive — ask upper years, everyone has their own favourite spot.")

# Slide 7 — Section header
slide = prs.slides.add_slide(L_SECTION)
slide.shapes.title.text = "MScAC Student Initiative Fund"

# Slide 8 — Past initiatives table
slide, body = add_title_content_slide("A few things the cohort has run this year", L_CONTENT)
layout_ph = L_CONTENT.placeholders[1]
left, top, width, height = layout_ph.left, layout_ph.top, layout_ph.width, layout_ph.height
sp = body._element
sp.getparent().remove(sp)
data = [
    ["Event", "What it was", "Funded"],
    ["Ski Trip (Blue Mountain)", "Full-day ski/snowboard trip, bus + lift tickets subsidized", "$2,760"],
    ["Bouldering & Social Mixer", "Climbing + dinner, peer mentorship for beginners", "$1,050"],
    ["Board Game Nights", "4-week recurring series", "$650"],
    ["CNY Celebration & Alumni Mixer", "Dinner, Lo Hei, karaoke, alumni invited back", "$2,315"],
    ["ML Trivia Night", "Dinner + trivia ahead of recruiting season", "$960"],
]
add_table(slide, len(data), 3, left, top, width, height, data,
          col_widths=[Inches(3.0), Inches(6.3), Inches(1.7)])
set_notes(slide, 'These span the full range of format — one-off outings, recurring series, cultural events, alumni-inclusive events. The fund isn\'t just "pizza and a room."')

# Slide 9 — Photo collage
from PIL import Image

slide = prs.slides.add_slide(L_TITLE_ONLY)
slide.shapes.title.text = "Look at all the fun we've had"
pics_dir = "../pictures"
cropped_dir = "cropped_pics"
os.makedirs(cropped_dir, exist_ok=True)
all_pics = sorted(os.listdir(pics_dir))
if len(all_pics) <= 8:
    pics = all_pics
else:
    step = len(all_pics) / 8
    pics = [all_pics[int(i * step)] for i in range(8)]
cols = 4
cell_w = Inches(3.2)
cell_h = Inches(2.35)
pad = Inches(0.1)
start_top = Inches(1.35)
start_left = Inches(0.2)

target_ratio = (cell_w - pad) / (cell_h - pad)

for i, pic in enumerate(pics):
    r, c = divmod(i, cols)
    left = start_left + c * cell_w
    top = start_top + r * cell_h

    src_path = os.path.join(pics_dir, pic)
    img = Image.open(src_path)
    w, h = img.size
    src_ratio = w / h

    if src_ratio > target_ratio:
        # source wider than target -> crop left/right
        new_w = int(h * target_ratio)
        offset = (w - new_w) // 2
        box = (offset, 0, offset + new_w, h)
    else:
        # source taller than target -> crop top/bottom
        new_h = int(w / target_ratio)
        offset = (h - new_h) // 2
        box = (0, offset, w, offset + new_h)

    cropped = img.crop(box)
    cropped_path = os.path.join(cropped_dir, f"crop_{i}.jpg")
    cropped.convert("RGB").save(cropped_path, quality=90)

    slide.shapes.add_picture(cropped_path, left, top, width=cell_w - pad, height=cell_h - pad)

# Slide 10 — Future ideas table
slide, body = add_title_content_slide("Other ideas to start you off this Fall", L_CONTENT)
layout_ph = L_CONTENT.placeholders[1]
left, top, width, height = layout_ph.left, layout_ph.top, layout_ph.width, layout_ph.height
sp = body._element
sp.getparent().remove(sp)
data = [
    ["Idea", "Description", "Cost"],
    ["K1 Speed Go-Karting (North York)", "Great fun near the Downsview Park TTC", "$82/person"],
    ["Basecamp Climbing (Queen West)", "Self-explanatory and we had dinner subsidised too", "$35/person"],
    ["Canada's Wonderland", "The GTA's rollercoaster park", "$45/person"],
    ["Activate Toronto (Stockyards)", "All-ages arcade park; 75-min active session", "$30/person"],
]
add_table(slide, len(data), 3, left, top, width, height, data,
          col_widths=[Inches(3.3), Inches(6.0), Inches(1.7)])
set_notes(slide, "These are priced per person so you can see the fund goes further than you'd think — none of these require a huge ask.")

# Slide 11 — Running an event (two content)
slide = prs.slides.add_slide(L_TWO)
slide.shapes.title.text = "Running an event"
left_ph, right_ph = [ph for ph in slide.placeholders if ph.placeholder_format.idx in (1, 2)]
set_bullets(left_ph.text_frame, [
    "In the office",
    ("Location: MScAC Office rooms like 9016", 1),
    ("Food: order delivery where possible — Mother's Dumpling (Chinese), Tahini's (shawarma), pizza (you'll get sick of this)", 1),
    ("Stock up on cutlery, plates and cups from Dollarama", 1),
    ("Buy big bottles of drinks from Shoppers", 1),
], font_size=16)
set_bullets(right_ph.text_frame, [
    "Outside the office",
    ("Show up, snap a picture or two, have fun", 1),
], font_size=16)

# Slide 12 — Rules / Guidelines
slide, body = add_title_content_slide("Guidelines from MScAC", L_CONTENT)
set_bullets(body.text_frame, [
    "Open to the entire cohort",
    "No alcohol, gambling / games of chance (stakes, betting pools, wagering)",
    "Proposal is one page: summary, goals, benefits, timeline, itemized budget (just use my Claude Code repo or something)",
    "You must be registered + in good academic standing (which is why us alumni need you to keep this up)",
], font_size=20)
source_box = slide.shapes.add_textbox(Inches(0.5), Inches(6.9), Inches(10), Inches(0.4))
tf = source_box.text_frame
tf.text = "(Source: MScAC Student-led Initiatives Guidelines)"
tf.paragraphs[0].font.size = Pt(12)
tf.paragraphs[0].font.italic = True
tf.paragraphs[0].font.color.rgb = RGBColor(0x77, 0x77, 0x77)

# Slide 13 — Timeline
slide, body = add_title_content_slide("Typical timeline", L_CONTENT)
set_bullets(body.text_frame, [
    "Submit a proposal 2 weeks ahead (the start of the semester is always a little busier)",
    "MScAC promises to confirm funding within 5 business days of your email",
    "After event: submit the attendance list (Partiful / written attendance sheet / Google Form), receipts, copy of original invite, and photos if applicable",
    "Expect reimbursement via cheque/direct deposit (if you're a TA) about 4-6 weeks after",
    "Past asks have ranged $650 – $2,760 depending on format and scale",
], font_size=18)

# Slide 14 — My workflow
slide, body = add_title_content_slide("Student Initiatives Process", L_CONTENT)
set_bullets(body.text_frame, [
    "1. Discord poll — gauge interest before writing anything",
    "2. Submit one-page proposal — real example at end of presentation",
    "3. Partiful — collect RSVPs + dietary restrictions once confirmed",
    "4. EVENT",
    "5. Submit receipts + attendee list + invitation",
    "6. Reimbursement lands in ~4-6 weeks",
], font_size=20)

# Slide 15 — Section header, transition to live doc walkthrough
slide = prs.slides.add_slide(L_SECTION)
slide.shapes.title.text = "What a real proposal actually looks like"
slide.placeholders[1].text = "(switching to the actual Bouldering proposal document)"

# Slide 16 — Section header
slide = prs.slides.add_slide(L_SECTION)
slide.shapes.title.text = "Want to run one yourself?"

# Slide 17 — Mentorship signup
slide, body = add_title_content_slide("Pick one, co-run it with me:", L_CONTENT)
set_bullets(body.text_frame, [
    "K1 Speed Go-Karting (round 2)",
    "Basecamp Climbing",
    "Canada's Wonderland",
    "Activate Toronto",
    "",
    "I'll help with the proposal, the budget, and the logistics to get you started.",
    "Sign up here: [Discord link — TBD]",
], font_size=20)

# Slide 18 — Closing
slide = prs.slides.add_slide(L_TITLE_ONLY)
slide.shapes.title.text = "Thanks!"
box = slide.shapes.add_textbox(Inches(0.5), Inches(2.0), Inches(6), Inches(1.5))
tf = box.text_frame
tf.word_wrap = True
tf.text = "[QR code to Discord signup — add once link is set]"
tf.paragraphs[0].font.size = Pt(20)

prs.save("student_life_session.pptx")
print("slide count:", len(prs.slides._sldIdLst))
